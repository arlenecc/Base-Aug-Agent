"""Comprehensive end-to-end RAG pipeline test.

Validates the ENTIRE pipeline from start to finish:
  ingestion -> text cleaning -> chunking -> vectorization ->
  semantic search -> reranking -> formatted output ->
  incremental ingest -> force reingest -> status ->
  tool integration -> agent flow -> clear.

The tests run in definition order within a single class so that the
stateful stages (ingest, re-ingest, clear) execute in the prescribed
sequence while sharing one ingested knowledge base.

Run:
  python3 -m pytest tests/test_rag_full_pipeline.py -v --tb=short
"""
from __future__ import annotations

import os

import pytest


# ---------------------------------------------------------------------------
# Test document factories
# ---------------------------------------------------------------------------

def _make_docx(path: str, title: str, paragraphs, table_rows=None) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=1)
    for p in paragraphs:
        doc.add_paragraph(p)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        table.style = "Table Grid"
        for i, row in enumerate(table_rows):
            for j, cell in enumerate(row):
                table.cell(i, j).text = cell
    doc.save(path)


def _make_xlsx(path: str, sheet_name: str, rows) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name
    for row in rows:
        ws.append(row)
    wb.save(path)


def _make_pptx(path: str, slides) -> None:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        for i, bullet in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
    prs.save(path)


def _create_test_documents(kb_dir: str) -> None:
    """Create 7 test files (.txt, .md, .csv, .html, .docx, .xlsx, .pptx).

    Each file carries meaningful Chinese + English content about a topic
    relevant to the search assertions below (RAG, 向量数据库, 知识库).
    """
    os.makedirs(kb_dir, exist_ok=True)

    # 1. .md — RAG技术
    with open(os.path.join(kb_dir, "rag_intro.md"), "w", encoding="utf-8") as f:
        f.write(
            "# RAG 技术介绍 (RAG Technology Overview)\n\n"
            "## 什么是 RAG / What is RAG\n\n"
            "RAG（Retrieval-Augmented Generation，检索增强生成）是一种结合信息检索与"
            "文本生成的技术架构。它在生成回答前先从外部知识库中检索相关文档片段，"
            "再将检索结果作为上下文提供给大语言模型，从而显著提高回答的准确性与时效性。\n\n"
            "## RAG 的核心组件 / Core Components\n\n"
            "1. **文档解析器（Document Parser）**：从 PDF、Word、Markdown 等格式中提取文本。\n"
            "2. **文本切片器（Text Chunker）**：将长文档分割成适合模型处理的短片段。\n"
            "3. **向量化引擎（Embedding Engine）**：将文本片段转换为向量表示。\n"
            "4. **向量数据库（Vector Store）**：存储和检索向量，常用 ChromaDB、FAISS。\n"
            "5. **重排序器（Reranker）**：对初步检索结果进行精确排序，提升相关性。\n\n"
            "## RAG 的优势 / Advantages\n\n"
            "相比传统的关键词搜索，RAG 使用语义向量搜索，能够理解查询的深层含义。"
            "通过 BGE Reranker 等模型进行重排序，可以进一步过滤不相关的结果，"
            "让最终答案更精确、更可信。\n"
        )

    # 2. .txt — 向量数据库
    with open(os.path.join(kb_dir, "vector_db.txt"), "w", encoding="utf-8") as f:
        f.write(
            "向量数据库对比分析 (Vector Database Comparison)\n\n"
            "向量数据库是专门用于存储和检索高维向量数据的数据库系统，"
            "在 AI 与 RAG 应用中扮演着关键角色。\n\n"
            "1. ChromaDB\n"
            "   - 开源、轻量级\n"
            "   - 支持多种嵌入模型\n"
            "   - 内置持久化存储\n"
            "   - 适合中小规模应用\n\n"
            "2. FAISS\n"
            "   - Facebook AI 开发\n"
            "   - 高性能向量检索\n"
            "   - 支持 GPU 加速\n"
            "   - 适合大规模检索场景\n\n"
            "3. Pinecone\n"
            "   - 云原生向量数据库\n"
            "   - 全托管服务\n"
            "   - 自动扩缩容\n"
            "   - 适合生产环境部署\n\n"
            "4. Weaviate\n"
            "   - 开源向量搜索引擎\n"
            "   - 支持 GraphQL 查询\n"
            "   - 内置多种向量化模块\n"
            "   - 适合复杂查询场景\n\n"
            "选择向量数据库时需要考虑：数据规模、查询延迟、部署复杂度、成本等因素。\n"
        )

    # 3. .csv — embedding models
    with open(os.path.join(kb_dir, "embedding_models.csv"), "w", encoding="utf-8") as f:
        f.write(
            "模型名称,维度,语言支持,特点\n"
            "all-MiniLM-L6-v2,384,multilingual,lightweight and fast\n"
            "bge-large-zh-v1.5,1024,Chinese,optimized for Chinese\n"
            "text-embedding-3-small,1536,multilingual,OpenAI official\n"
            "m3e-base,768,Chinese+English,popular in Chinese community\n"
            "bge-m3,1024,multilingual,multilingual retrieval\n"
        )

    # 4. .html — 知识库构建 (with script/style/nav/header/footer to test cleaning)
    with open(os.path.join(kb_dir, "kb_construction.html"), "w", encoding="utf-8") as f:
        f.write(
            "<!DOCTYPE html>\n"
            "<html>\n"
            "<head><title>知识库构建最佳实践</title></head>\n"
            "<body>\n"
            "<header>导航栏 (should be cleaned)</header>\n"
            "<nav>菜单 nav (should be cleaned)</nav>\n"
            "<main>\n"
            "<h1>知识库构建最佳实践 (Knowledge Base Best Practices)</h1>\n"
            "<p>构建高质量的知识库是 RAG 系统成功的关键。</p>\n"
            "<h2>文档质量控制</h2>\n"
            "<p>确保知识库中的文档内容准确、完整、格式规范。"
            "避免包含大量扫描件或图片型 PDF，优先使用文本型文档。</p>\n"
            "<h2>文本切片策略</h2>\n"
            "<p>切片大小应根据模型上下文窗口合理设置。一般建议 500 token 一个切片，"
            "保留约 10% 的重叠区域，避免信息在切片边界处丢失。</p>\n"
            "<h2>元数据管理</h2>\n"
            "<p>为每个文档添加标题、来源、日期等元数据信息，方便检索时过滤和溯源。</p>\n"
            "<h2>定期更新</h2>\n"
            "<p>知识库需要定期更新，确保信息的时效性。建议设置自动同步机制。</p>\n"
            "</main>\n"
            "<footer>页脚 footer (should be cleaned)</footer>\n"
            "<script>console.log('script should be removed');</script>\n"
            "<style>body { color: red; }</style>\n"
            "</body>\n"
            "</html>\n"
        )

    # 5. .docx — RAG架构详解
    _make_docx(
        os.path.join(kb_dir, "rag_architecture.docx"),
        title="RAG 架构详解 (RAG Architecture Deep Dive)",
        paragraphs=[
            "RAG（检索增强生成）架构由三个核心阶段组成：检索、增强、生成。"
            "Retrieval-Augmented Generation consists of retrieval, augmentation, "
            "and generation stages.",
            "在检索阶段，系统将用户查询向量化，并在向量数据库中搜索最相关的文档片段。"
            "During retrieval, the user query is embedded and matched against a "
            "vector database of document chunks.",
            "在增强阶段，检索到的文档片段与原始查询拼接，形成增强后的上下文。"
            "In augmentation, retrieved chunks are concatenated with the query to "
            "form an enriched context for the LLM.",
            "在生成阶段，大语言模型基于增强上下文生成最终回答，"
            "从而减少幻觉并提高事实准确性。In generation, the LLM produces the "
            "final answer grounded in the retrieved context, reducing hallucinations.",
            "RAG 系统通常结合重排序器（如 BGE Reranker）对检索结果精排，"
            "只保留最相关的片段送入生成模型，以节省上下文窗口。",
        ],
        table_rows=[
            ["阶段 Stage", "作用 Role", "典型组件 Components"],
            ["检索 Retrieval", "召回相关文档", "Embedding + Vector Store"],
            ["增强 Augmentation", "拼接上下文", "Prompt Builder"],
            ["生成 Generation", "生成回答", "LLM"],
        ],
    )

    # 6. .xlsx — 向量数据库对比
    _make_xlsx(
        os.path.join(kb_dir, "vector_db_comparison.xlsx"),
        sheet_name="VectorDBs",
        rows=[
            ["名称 Name", "类型 Type", "维度支持 Dim", "特点 Features"],
            ["ChromaDB", "开源 Open-source", "任意 Any", "轻量、内置持久化 lightweight"],
            ["FAISS", "开源 Open-source", "任意 Any", "GPU 加速、高性能 GPU accelerated"],
            ["Pinecone", "云服务 Cloud", "任意 Any", "全托管 fully managed"],
            ["Weaviate", "开源 Open-source", "任意 Any", "GraphQL、模块化 modular"],
            ["Milvus", "开源 Open-source", "任意 Any", "分布式、大规模 distributed"],
        ],
    )

    # 7. .pptx — 知识库构建流程
    _make_pptx(
        os.path.join(kb_dir, "kb_process.pptx"),
        slides=[
            (
                "知识库构建流程 (Knowledge Base Construction)",
                [
                    "第一步：收集文档 Collect documents (PDF, Word, Markdown).",
                    "第二步：解析提取 Extract text with document parsers.",
                    "第三步：清洗文本 Clean and normalize text.",
                    "第四步：切片 Chunk text into token-sized pieces.",
                ],
            ),
            (
                "向量化与检索 (Embedding and Retrieval)",
                [
                    "将切片向量化 Embed chunks with an embedding model.",
                    "存入向量数据库 Store vectors in a vector database.",
                    "查询时检索最相关片段 Retrieve top-k chunks at query time.",
                    "使用重排序器精排 Rerank candidates for precision.",
                ],
            ),
            (
                "RAG 集成 (RAG Integration)",
                [
                    "检索结果作为上下文注入提示 Inject retrieved chunks into prompt.",
                    "大语言模型生成最终回答 LLM generates grounded answer.",
                    "知识库可定期更新 KB can be updated incrementally.",
                ],
            ),
        ],
    )


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(scope="class")
def kb_dir(tmp_path_factory):
    d = tmp_path_factory.mktemp("kb")
    _create_test_documents(str(d))
    return str(d)


@pytest.fixture(scope="class")
def workspace(tmp_path_factory):
    return str(tmp_path_factory.mktemp("ws"))


@pytest.fixture(scope="class")
def engine(workspace, kb_dir, rag_engine_factory):
    """A shared RAGEngine. Ingestion happens in the first test."""
    return rag_engine_factory(workspace=workspace, knowledge_base=kb_dir)


@pytest.fixture(scope="class")
def shared():
    """Cross-test shared state (e.g. initial chunk count for consistency checks)."""
    return {}


# ---------------------------------------------------------------------------
# Ordered end-to-end tests
# ---------------------------------------------------------------------------

class TestRagFullPipeline:
    """13 stages executed in definition order over a shared ingested KB."""

    # -- 1. ingest --------------------------------------------------------

    def test_full_pipeline_ingest(self, engine, kb_dir):
        stats = engine.ingest(force=True)

        assert "error" not in stats, f"ingest returned error: {stats.get('error')}"
        assert stats["files_found"] >= 7, (
            f"expected >=7 files found, got {stats['files_found']}"
        )
        assert (
            stats["files_extracted"] + stats["files_skipped"] >= 7
        ), f"not enough files processed: {stats}"
        assert stats["chunks"] > 0, "no chunks generated"
        assert len(stats["errors"]) <= 1, (
            f"too many extraction errors: {stats['errors']}"
        )

        # markdown cache must exist
        assert os.path.isdir(engine.markdown_dir), "markdown dir not created"
        md_files = [f for f in os.listdir(engine.markdown_dir) if f.endswith(".md")]
        assert len(md_files) >= 7, f"expected >=7 markdown files, got {len(md_files)}"

    # -- 2. markdown content quality --------------------------------------

    def test_markdown_content_quality(self, engine):
        md_dir = engine.markdown_dir
        assert os.path.isdir(md_dir)
        md_files = [f for f in os.listdir(md_dir) if f.endswith(".md")]
        assert len(md_files) >= 7

        for md_name in md_files:
            with open(os.path.join(md_dir, md_name), "r", encoding="utf-8") as f:
                content = f.read()

            assert content.strip(), f"empty markdown: {md_name}"
            # HTML tags must be cleaned out (especially from the .html-derived file)
            assert "<script>" not in content.lower(), f"script tag in {md_name}"
            assert "<style>" not in content.lower(), f"style tag in {md_name}"
            assert "<html>" not in content.lower(), f"html tag in {md_name}"
            # whitespace must be normalized (no 3+ consecutive newlines)
            assert "\n\n\n" not in content, f"3+ consecutive newlines in {md_name}"

    # -- 3. chunking correctness ------------------------------------------

    def test_chunking_correctness(self, kb_dir):
        from src.agent.rag.parsers import extract_directory
        from src.agent.rag.cleaner import clean_text, normalize_markdown
        from src.agent.rag.chunker import chunk_documents

        extracted = extract_directory(kb_dir, recursive=True)
        assert len(extracted) >= 7, f"only extracted {len(extracted)} files"

        documents = []
        combined_source = ""
        for filepath, raw in extracted:
            cleaned = normalize_markdown(clean_text(raw))
            if cleaned.strip():
                documents.append({"source": filepath, "text": cleaned})
                combined_source += cleaned + "\n"

        chunks = chunk_documents(documents)

        assert len(chunks) > 0, "no chunks produced"
        # more files should yield a reasonable number of chunks
        assert len(chunks) >= len(documents), (
            f"chunk count {len(chunks)} < document count {len(documents)}"
        )

        chunk_texts = [c["text"] for c in chunks]
        # each chunk non-empty
        assert all(t.strip() for t in chunk_texts), "found empty chunk"
        # no duplicate chunks
        assert len(chunk_texts) == len(set(chunk_texts)), "duplicate chunks found"
        # chunk text must come from the source documents
        for c in chunks:
            assert c["text"] in combined_source, (
                "chunk text not found in source documents"
            )

    # -- 4. vector storage persistence ------------------------------------

    def test_vector_storage_persistence(self, engine, workspace, kb_dir, shared, rag_engine_factory):
        store = engine._get_store()
        count_before = store.count()
        assert count_before > 0, "vector store is empty after ingest"

        sources = store.list_sources()
        assert len(sources) >= 7, f"only {len(sources)} sources in store"
        # every source file should be represented (compare by basename)
        source_basenames = {os.path.basename(s) for s in sources}
        expected = {
            "rag_intro.md", "vector_db.txt", "embedding_models.csv",
            "kb_construction.html", "rag_architecture.docx",
            "vector_db_comparison.xlsx", "kb_process.pptx",
        }
        missing = expected - source_basenames
        assert not missing, f"missing sources in store: {missing}"

        # brand-new engine pointing at the same workspace must see the vectors
        new_engine = rag_engine_factory(workspace=workspace, knowledge_base=kb_dir)
        new_store = new_engine._get_store()
        assert new_store.count() == count_before, (
            "persistence broken: new engine sees different vector count"
        )
        shared["persisted_count"] = count_before

    # -- 5. search returns relevant results -------------------------------

    def test_search_returns_relevant_results(self, engine):
        cases = [
            ("RAG", "RAG"),
            ("向量数据库", None),  # vector database
            ("知识库", None),       # knowledge base
        ]
        for query, must_contain in cases:
            results = engine.search(query, top_k=3)
            assert len(results) <= 3, f"top_k exceeded for query {query!r}"
            assert len(results) >= 1, f"no results for query {query!r}"
            for r in results:
                assert "text" in r, "result missing 'text'"
                assert "source" in r, "result missing 'source'"
                assert "score" in r, "result missing 'score'"

            combined = " ".join(r["text"] for r in results)
            if query == "RAG":
                assert "RAG" in combined, f"RAG query did not return RAG content: {combined[:120]}"
            elif query == "向量数据库":
                assert ("向量数据库" in combined or "ChromaDB" in combined
                        or "FAISS" in combined), (
                    f"向量数据库 query returned irrelevant content: {combined[:120]}"
                )
            elif query == "知识库":
                assert ("知识库" in combined or "knowledge" in combined.lower()), (
                    f"知识库 query returned irrelevant content: {combined[:120]}"
                )

    # -- 6. reranker scoring ----------------------------------------------

    def test_reranker_scoring(self, engine):
        results = engine.search("RAG 架构与向量数据库", top_k=3)
        assert len(results) >= 1, "no results for reranker test"

        scores = [r["score"] for r in results]
        # results must be sorted by score descending (most relevant first)
        assert scores == sorted(scores, reverse=True), (
            f"results not sorted by score desc: {scores}"
        )
        # scores must be in valid range [0, 1]
        for s in scores:
            assert 0.0 <= s <= 1.0, f"score out of [0,1] range: {s}"

    # -- 7. formatted output ----------------------------------------------

    def test_search_formatted_output(self, engine):
        output = engine.search_formatted("什么是RAG技术")
        assert isinstance(output, str), "formatted output must be a string"
        assert output.strip(), "formatted output is empty"
        assert "来源:" in output, "formatted output missing '来源:' markers"
        assert "相关度:" in output, "formatted output missing '相关度:' markers"

    # -- 8. incremental ingest (skip cached) ------------------------------

    def test_incremental_ingest_skip(self, engine, shared):
        store = engine._get_store()
        count_before = store.count()

        stats = engine.ingest(force=False)
        assert "error" not in stats
        assert stats["files_skipped"] == stats["files_found"], (
            f"expected all files skipped, got found={stats['files_found']} "
            f"skipped={stats['files_skipped']}"
        )
        assert stats["files_extracted"] == 0, (
            f"force=False should not re-extract, got {stats['files_extracted']}"
        )
        # no new chunks added: count unchanged
        assert store.count() == count_before, (
            "incremental ingest changed vector count unexpectedly"
        )

    # -- 9. force reingest ------------------------------------------------

    def test_force_reingest(self, engine, shared):
        initial = shared.get("persisted_count")
        assert initial and initial > 0, "prior ingest did not persist chunks"

        stats = engine.ingest(force=True)
        assert "error" not in stats
        assert stats["files_extracted"] == stats["files_found"], (
            f"force=True should reprocess all files, got "
            f"extracted={stats['files_extracted']} found={stats['files_found']}"
        )
        assert stats["files_skipped"] == 0, "force=True should skip nothing"
        # chunk count must be consistent with the previous full ingest
        assert stats["chunks"] == initial, (
            f"chunk count inconsistent: {stats['chunks']} vs {initial}"
        )
        assert engine._get_store().count() == initial, (
            "store count changed after force reingest"
        )

    # -- 10. rag status ---------------------------------------------------

    def test_rag_status(self, engine):
        status = engine.status()
        assert status["chunks_stored"] > 0, "status reports no chunks"
        assert len(status["sources"]) >= 7, "status sources list is empty"
        assert status["has_knowledge_base"] is True, "has_knowledge_base flag wrong"

    # -- 11. rag tool integration -----------------------------------------

    def test_rag_tool_integration(self, kb_dir, tmp_path, rag_engine_factory):
        from src.agent.tools.rag_tool import RagSearchTool

        eng = rag_engine_factory(workspace=str(tmp_path), knowledge_base=kb_dir)
        stats = eng.ingest(force=True)
        assert stats["chunks"] > 0, "tool-integration engine failed to ingest"

        tool = RagSearchTool(eng)
        result = tool.run(query="RAG")
        assert result.success is True, f"tool run failed: {result.error}"
        assert result.output, "tool returned empty output"
        assert "来源:" in result.output, "tool output missing '来源:'"

        # graceful handling of empty / missing query
        empty_result = tool.run(query="")
        assert empty_result.success is False, "empty query should fail gracefully"
        assert empty_result.error, "empty query should report an error"

    # -- 12. agent RAG-enhanced response ----------------------------------

    def test_agent_rag_enhanced_response(self, engine):
        """Simulate the agent tool-call flow and verify RAG context is usable."""
        from src.agent.tools.rag_tool import RagSearchTool

        # 1. User asks a question
        user_query = "什么是RAG技术"

        # 2. Agent decides to call the rag_search tool
        tool = RagSearchTool(engine)
        tool_result = tool.run(query=user_query)
        assert tool_result.success is True, (
            f"rag_search tool failed: {tool_result.error}"
        )

        # 3. Tool returns formatted results that would be injected into history
        context = tool_result.output
        assert isinstance(context, str) and context.strip(), (
            "tool returned empty context for agent"
        )

        # 4. The context must carry enough information to answer the question
        relevant_keywords = ["RAG", "检索增强生成", "Retrieval"]
        matched = [kw for kw in relevant_keywords if kw in context]
        assert matched, (
            f"RAG context lacks relevant keywords; got: {context[:200]!r}"
        )
        assert "来源:" in context, "context missing source attribution"
        assert "相关度:" in context, "context missing relevance markers"

    # -- 13. clear and verify empty ---------------------------------------

    def test_clear_and_verify_empty(self, engine):
        engine.clear()

        store = engine._get_store()
        assert store.count() == 0, "store not empty after clear()"

        assert engine.search("RAG") == [], "search should return empty after clear"

        md_dir = engine.markdown_dir
        if os.path.isdir(md_dir):
            leftovers = [f for f in os.listdir(md_dir) if f.endswith(".md")]
            assert leftovers == [], f"markdown dir not cleaned: {leftovers}"
