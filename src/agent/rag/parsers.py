"""Document parsers for Word, Excel, PPT, PDF, EPUB, MOBI, AZW3 and more.

Each parser accepts a file path and returns extracted text. Lazy imports keep
optional dependencies optional — only the formats you actually use need to be
installed.
"""
from __future__ import annotations

import logging
import os
import re
import tempfile
import threading
from pathlib import Path
from typing import Iterator, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 全局 OCR 引擎（单例）+ 并发信号量
# ---------------------------------------------------------------------------
# RapidOCR ONNX Runtime 单实例约 500MB 内存。如果每个 worker 各加载一份，
# 4 worker 同时遇到图片 PDF 时会占用 ~2GB 仅仅用于 OCR 模型。
# 改为全局单例：所有 worker 共享同一个引擎实例，通过信号量串行执行
# （RapidOCR 可能非线程安全，串行调用 engine(img) 避免竞态）。
_OCR_SEMAPHORE = threading.Semaphore(1)
_OCR_ENGINE = None  # 全局单例；None 表示尚未尝试加载
_OCR_ENGINE_LOCK = threading.Lock()

# ---------------------------------------------------------------------------
# Format → extension mapping
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx", ".ppt",
    ".pdf",
    ".epub",
    ".mobi", ".azw3", ".azw",
    ".txt", ".md", ".markdown", ".rst",
    ".csv", ".tsv",
    ".html", ".htm",
}

# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def extract_text(filepath: str) -> str:
    """Extract text from a document. Returns the full text as a string."""
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    suffix = path.suffix.lower()

    # 文件类型→解析器映射（用于日志）
    _PARSER_LABELS = {
        ".docx": "Word 文档", ".doc": "Word 文档",
        ".xlsx": "Excel 表格", ".xls": "Excel 表格",
        ".pptx": "PPT 幻灯片", ".ppt": "PPT 幻灯片",
        ".pdf": "PDF 文档",
        ".epub": "EPUB 电子书",
        ".mobi": "MOBI 电子书", ".azw3": "AZW3 电子书", ".azw": "AZW 电子书",
        ".txt": "纯文本", ".md": "Markdown", ".markdown": "Markdown", ".rst": "reStructuredText",
        ".csv": "CSV 表格", ".tsv": "TSV 表格",
        ".html": "HTML 网页", ".htm": "HTML 网页",
    }
    label = _PARSER_LABELS.get(suffix, suffix)
    logger.debug("    解析文件类型: %s (%s)", path.name, label)

    if suffix in (".docx", ".doc"):
        return _extract_docx(path)
    elif suffix in (".xlsx", ".xls"):
        return _extract_xlsx(path)
    elif suffix in (".pptx", ".ppt"):
        return _extract_pptx(path)
    elif suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix == ".epub":
        return _extract_epub(path)
    elif suffix in (".mobi", ".azw3", ".azw"):
        return _extract_mobi(path)
    elif suffix in (".txt", ".md", ".markdown", ".rst"):
        return _extract_text(path)
    elif suffix in (".csv", ".tsv"):
        return _extract_csv(path)
    elif suffix in (".html", ".htm"):
        return _extract_html(path)
    else:
        raise ValueError(f"Unsupported format: {suffix}")


# ---------------------------------------------------------------------------
# Individual parsers
# ---------------------------------------------------------------------------


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required for Word files. pip install python-docx")

    doc = Document(str(path))
    parts: List[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is required for Excel files. pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts: List[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"# Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PPT files. pip install python-pptx")

    prs = Presentation(str(path))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))
        if slide_parts:
            parts.append(f"## Slide {i}\n\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF. Falls back to OCR for image-based pages."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF is required for PDF files. pip install PyMuPDF")

    doc = fitz.open(str(path))
    parts: List[str] = []
    ocr_engine = None       # lazy-init, 复用同一个引擎实例
    ocr_unavailable = False  # OCR 引擎加载失败标记，避免每页重复尝试
    ocr_needed = False

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()
        if text:
            # 文字版页面：直接使用提取的文本，不用 OCR
            parts.append(text)
            continue

        # 无文字层 → 可能是图片版页面。检查是否含图片
        if not page.get_images(full=True):
            # 既无文字也无图片 → 跳过（可能是空白页）
            continue

        # 图片版页面 → 走 OCR（用全局信号量限制并发，避免多 worker 同时跑 OCR 导致 OOM）
        ocr_needed = True
        if ocr_unavailable:
            continue
        if ocr_engine is None:
            ocr_engine = _get_ocr_engine()
            if ocr_engine is None:
                ocr_unavailable = True
                continue
        with _OCR_SEMAPHORE:
            # 再次检查引擎：可能在等信号量期间其它 worker 已判定不可用
            if ocr_engine is None:
                ocr_unavailable = True
                continue
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            ocr_text = _ocr_image_cached(ocr_engine, img_bytes, page_num + 1, path.name)
        if ocr_text:
            parts.append(ocr_text)

    doc.close()

    if not parts:
        return ""
    if ocr_needed:
        logger.info("PDF '%s': used OCR on some pages", path.name)
    return "\n\n".join(parts)


def _get_ocr_engine():
    """Lazy-load RapidOCR engine as a global singleton.

    Returns None if not installed. All workers share the same engine instance
    (~500MB) — without this, 4 workers loading 4 instances would OOM. The
    _OCR_SEMAPHORE serializes actual OCR calls so the shared engine is only
    invoked by one thread at a time.
    """
    global _OCR_ENGINE
    if _OCR_ENGINE is not None:
        return _OCR_ENGINE
    with _OCR_ENGINE_LOCK:
        if _OCR_ENGINE is not None:
            return _OCR_ENGINE
        try:
            from rapidocr_onnxruntime import RapidOCR
            _OCR_ENGINE = RapidOCR()
            logger.info("RapidOCR engine loaded (global singleton)")
            return _OCR_ENGINE
        except ImportError:
            logger.warning(
                "rapidocr-onnxruntime is required for OCR on image-based PDFs. "
                "pip install rapidocr-onnxruntime"
            )
            return None
        except Exception as e:
            logger.warning("Failed to init RapidOCR: %s", e)
            return None


def _ocr_image_cached(engine, img_bytes: bytes, page_num: int, filename: str) -> str:
    """Run OCR using a pre-loaded engine instance."""
    try:
        result, _ = engine(img_bytes)
        if not result:
            return ""
        lines = []
        for item in result:
            text = item[1] if isinstance(item, (list, tuple)) and len(item) > 1 else str(item)
            if text.strip():
                lines.append(text)
        return "\n".join(lines)
    except Exception as e:
        logger.warning("OCR failed on page %d of '%s': %s", page_num, filename, e)
        return ""


def _extract_epub(path: Path) -> str:
    try:
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("ebooklib and beautifulsoup4 are required for EPUB. pip install ebooklib beautifulsoup4")

    book = epub.read_epub(str(path))
    parts: List[str] = []
    for item in book.get_items_of_type(9):  # ITEM_DOCUMENT = 9
        content = item.get_content().decode("utf-8", errors="replace")
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text(separator="\n")
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_mobi(path: Path) -> str:
    """Extract text from MOBI/AZW3 files.

    Strategy: try to convert via calibre ebook-convert CLI first (most reliable),
    then fall back to a basic raw-text extraction.
    """
    # Try calibre's ebook-convert first
    import subprocess
    import shutil

    if shutil.which("ebook-convert"):
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
            txt_path = tmp.name
        try:
            result = subprocess.run(
                ["ebook-convert", str(path), txt_path],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0 and os.path.exists(txt_path):
                with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read()
        except (subprocess.TimeoutExpired, Exception) as e:
            logger.warning("ebook-convert failed for '%s': %s", path.name, e)
        finally:
            try:
                os.unlink(txt_path)
            except Exception:
                pass

    # Fallback: basic raw extraction
    try:
        return _extract_mobi_raw(path)
    except Exception as e:
        logger.warning("Raw MOBI extraction failed for '%s': %s", path.name, e)
        return ""


def _extract_mobi_raw(path: Path) -> str:
    """Basic MOBI text extraction without external tools."""
    with open(path, "rb") as f:
        data = f.read()

    # Try to find text content in the raw bytes
    text_parts: List[str] = []
    # Look for UTF-8 text segments
    try:
        decoded = data.decode("utf-8", errors="ignore")
        # Extract readable text blocks (sequences of printable chars)
        blocks = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\w\s.,!?;:()\[\]{}\'\"\-–—]{30,}', decoded)
        text_parts.extend(blocks)
    except Exception:
        pass

    # Also try Latin-1
    if not text_parts:
        try:
            decoded = data.decode("latin-1", errors="ignore")
            blocks = re.findall(r'[\w\s.,!?;:()\[\]{}\'\"\-–—]{30,}', decoded)
            text_parts.extend(blocks)
        except Exception:
            pass

    return "\n\n".join(text_parts)


def _extract_text(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_csv(path: Path) -> str:
    import csv
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        sample = f.read(4096)
        f.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel  # default: comma-delimited
        reader = csv.reader(f, dialect)
        rows = list(reader)
    if not rows:
        return ""
    lines = [" | ".join(r) for r in rows]
    return "\n".join(lines)


def _extract_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        # Fallback: basic regex stripping
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        content = re.sub(r'<[^>]+>', ' ', content)
        content = re.sub(r'\s+', ' ', content)
        return content.strip()

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # Remove script/style tags
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n")


# ---------------------------------------------------------------------------
# Directory-level extraction
# ---------------------------------------------------------------------------


def extract_directory(
    dirpath: str,
    recursive: bool = True,
    extensions: Optional[set] = None,
    errors: Optional[List[str]] = None,
) -> List[tuple]:
    """Walk a directory and extract text from all supported files.

    Returns a list of (filepath, extracted_text) tuples.
    If *errors* is provided, extraction failure messages are appended to it.
    """
    return list(iter_directory(dirpath, recursive=recursive,
                               extensions=extensions, errors=errors))


def iter_directory(
    dirpath: str,
    recursive: bool = True,
    extensions: Optional[set] = None,
    errors: Optional[List[str]] = None,
) -> Iterator[Tuple[str, str]]:
    """流式版本的 extract_directory：逐文件 yield (filepath, text)。

    避免大知识库一次性把全部文件全文累积在内存（H1 隐患）。
    调用方消费完一个文件后即可释放该文件的全文字符串。
    """
    if extensions is None:
        extensions = SUPPORTED_EXTENSIONS

    root = Path(dirpath)
    iterator = root.rglob("*") if recursive else root.glob("*")
    for filepath in iterator:
        if not filepath.is_file():
            continue
        if filepath.suffix.lower() not in extensions:
            continue
        # Skip hidden files
        if filepath.name.startswith("."):
            continue
        try:
            text = extract_text(str(filepath))
            if text.strip():
                yield (str(filepath), text)
                logger.info("Extracted: %s (%d chars)", filepath.name, len(text))
            else:
                msg = f"{filepath.name}: extracted text is empty"
                logger.warning(msg)
                if errors is not None:
                    errors.append(msg)
        except Exception as e:
            msg = f"{filepath.name}: {e}"
            logger.warning("Failed to extract '%s': %s", filepath.name, e)
            if errors is not None:
                errors.append(msg)